"""File processing service for extracting text from various document formats."""

import asyncio
import io
import logging
import shutil
import subprocess
import tempfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image

logger = logging.getLogger(__name__)

# Thread pool for CPU-intensive operations
_executor = ThreadPoolExecutor(max_workers=2)

# Cache for poppler path
_poppler_path: Optional[str] = None


def _find_poppler_path() -> Optional[str]:
    """Find poppler binaries path (pdftoppm)."""
    # Common locations
    possible_paths = [
        "/usr/bin",  # Linux standard
        "/usr/local/bin",  # macOS homebrew
        "/opt/homebrew/bin",  # macOS Apple Silicon
    ]

    for path in possible_paths:
        pdftoppm = Path(path) / "pdftoppm"
        if pdftoppm.exists():
            logger.info("[FILE_PROCESSOR] Found poppler at: %s", path)
            return path

    # Try to find via which/where
    result = shutil.which("pdftoppm")
    if result:
        path = str(Path(result).parent)
        logger.info("[FILE_PROCESSOR] Found poppler via which: %s", path)
        return path

    logger.warning("[FILE_PROCESSOR] poppler (pdftoppm) not found in common paths")
    return None


def get_poppler_path() -> Optional[str]:
    """Get cached poppler path."""
    global _poppler_path
    if _poppler_path is None:
        _poppler_path = _find_poppler_path()
    return _poppler_path


def verify_poppler_installed() -> bool:
    """Verify that poppler is installed and working."""
    path = get_poppler_path()
    if not path:
        return False

    pdftoppm = Path(path) / "pdftoppm"
    if not pdftoppm.exists():
        return False

    try:
        result = subprocess.run(
            [str(pdftoppm), "-v"],
            capture_output=True,
            text=True,
            timeout=5
        )
        if result.returncode == 0:
            version = result.stderr.strip() or result.stdout.strip()
            logger.info("[FILE_PROCESSOR] Poppler verified: %s", version.split('\n')[0])
            return True
    except Exception as e:
        logger.warning("[FILE_PROCESSOR] Poppler verification failed: %s", e)

    return False


class FileProcessor:
    """Extract text content from various file formats."""

    SUPPORTED_IMAGE_TYPES = {
        'image/jpeg', 'image/png', 'image/gif', 'image/webp', 'image/bmp'
    }
    SUPPORTED_DOCUMENT_TYPES = {
        'application/pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # docx
        'application/msword',  # doc
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',  # xlsx
        'application/vnd.ms-excel',  # xls
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # pptx
        'application/vnd.ms-powerpoint',  # ppt
        'text/plain',
        'text/markdown',
        'text/csv',
    }
    # Types that require Vision AI analysis (converted to images)
    VISION_TYPES = {
        'application/pdf',
        *SUPPORTED_IMAGE_TYPES,
    }

    @classmethod
    def is_supported(cls, mime_type: str) -> bool:
        """Check if file type is supported."""
        return mime_type in cls.SUPPORTED_IMAGE_TYPES or mime_type in cls.SUPPORTED_DOCUMENT_TYPES

    @classmethod
    async def extract_text(
        cls,
        file_content: bytes,
        mime_type: str,
        file_name: Optional[str] = None
    ) -> Tuple[str, dict]:
        """
        Extract text from file content.

        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        logger.info("[FILE_PROCESSOR] Starting text extraction for: %s (MIME: %s, Size: %d bytes)",
                    file_name, mime_type, len(file_content))

        metadata = {
            'original_filename': file_name,
            'mime_type': mime_type,
            'file_size': len(file_content),
        }

        try:
            if mime_type in cls.SUPPORTED_IMAGE_TYPES:
                logger.info("[FILE_PROCESSOR] Processing as IMAGE")
                text = await cls._extract_from_image(file_content, mime_type)
                metadata['content_type'] = 'image'
            elif mime_type == 'application/pdf':
                logger.info("[FILE_PROCESSOR] Processing as PDF")
                text = await cls._extract_from_pdf(file_content)
                metadata['content_type'] = 'pdf'
            elif mime_type in (
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'application/msword'
            ):
                logger.info("[FILE_PROCESSOR] Processing as WORD")
                text = await cls._extract_from_word(file_content)
                metadata['content_type'] = 'word'
            elif mime_type in (
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'application/vnd.ms-excel'
            ):
                logger.info("[FILE_PROCESSOR] Processing as EXCEL")
                text = await cls._extract_from_excel(file_content)
                metadata['content_type'] = 'excel'
            elif mime_type in (
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'application/vnd.ms-powerpoint'
            ):
                logger.info("[FILE_PROCESSOR] Processing as POWERPOINT")
                text = await cls._extract_from_powerpoint(file_content)
                metadata['content_type'] = 'powerpoint'
            elif mime_type in ('text/plain', 'text/markdown', 'text/csv'):
                logger.info("[FILE_PROCESSOR] Processing as TEXT")
                text = file_content.decode('utf-8', errors='ignore')
                metadata['content_type'] = 'text'
            else:
                logger.error("[FILE_PROCESSOR] Unsupported MIME type: %s", mime_type)
                raise ValueError(f"Unsupported file type: {mime_type}")

            metadata['extracted_length'] = len(text)
            logger.info("[FILE_PROCESSOR] Extraction complete. Extracted %d characters", len(text))
            return text, metadata

        except Exception as e:
            logger.exception("[FILE_PROCESSOR] Failed to extract text from file: %s - Error: %s", file_name, e)
            raise ValueError(f"Failed to process file: {str(e)}") from e

    @classmethod
    async def _extract_from_image(cls, content: bytes, mime_type: str) -> str:
        """Store image info for later analysis via MCP GLM Vision.

        Note: We don't use OCR here. Use the 'analyze_image' tool via MCP
        for AI-powered image analysis with GLM Vision.
        """
        logger.info("[FILE_PROCESSOR] Processing image (%s, %d bytes) - use analyze_image tool for AI analysis", mime_type, len(content))
        try:
            from PIL import Image
            image = Image.open(io.BytesIO(content))
            width, height = image.size
            mode = image.mode
            logger.info("[FILE_PROCESSOR] Image info: %dx%d, mode=%s", width, height, mode)
            return f"[Image: {width}x{height} pixels, {mime_type} - Use 'analyze_image' tool for AI-powered analysis]"
        except Exception as e:
            logger.warning("[FILE_PROCESSOR] Could not read image info: %s", e)
            return f"[Image file: {mime_type} - Use 'analyze_image' tool for AI-powered analysis]"

    @classmethod
    async def _extract_from_pdf(cls, content: bytes) -> str:
        """Extract text from PDF using pypdf (native text extraction).

        This is the primary method for PDF text extraction - fast and reliable
        for PDFs with embedded text. Falls back to Vision API in knowledge.py
        for scanned PDFs.
        """
        return await cls.extract_pdf_text(content)

    @classmethod
    async def extract_pdf_text(cls, content: bytes) -> Tuple[str, dict]:
        """
        Extract text from PDF using pypdf (native text extraction).

        This is fast and reliable for PDFs with embedded text.
        For scanned PDFs, the caller should use Vision API as fallback.

        Args:
            content: PDF file content as bytes

        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        logger.info("[FILE_PROCESSOR] Starting PDF text extraction with pypdf")
        metadata = {
            'extraction_method': 'pypdf',
            'pages_processed': 0,
            'pages_with_text': 0,
        }

        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(content))
            total_pages = len(reader.pages)
            metadata['total_pages'] = total_pages

            logger.info("[FILE_PROCESSOR] PDF has %d pages", total_pages)

            text_parts = []
            pages_with_text = 0

            for i, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(f"--- Page {i + 1} ---")
                        text_parts.append(page_text.strip())
                        pages_with_text += 1
                except Exception as e:
                    logger.warning("[FILE_PROCESSOR] Failed to extract text from page %d: %s", i + 1, e)

            metadata['pages_processed'] = total_pages
            metadata['pages_with_text'] = pages_with_text

            extracted_text = '\n\n'.join(text_parts)

            # Check if we got meaningful text
            if not extracted_text.strip() or pages_with_text == 0:
                logger.warning("[FILE_PROCESSOR] PDF text extraction returned empty result - likely scanned PDF")
                metadata['is_scanned'] = True
                return "", metadata

            logger.info("[FILE_PROCESSOR] PDF text extraction complete: %d chars from %d/%d pages",
                       len(extracted_text), pages_with_text, total_pages)

            return extracted_text, metadata

        except ImportError:
            logger.error("[FILE_PROCESSOR] pypdf not installed. Install with: pip install pypdf")
            metadata['error'] = 'pypdf not installed'
            return "", metadata
        except Exception as e:
            logger.exception("[FILE_PROCESSOR] PDF text extraction failed: %s", e)
            metadata['error'] = str(e)
            return "", metadata

    @classmethod
    async def extract_word_text(cls, content: bytes) -> Tuple[str, dict]:
        """
        Extract text from Word document (.docx only).

        Note: .doc (Word 97-2003) format is NOT supported.
        Users must convert .doc to .docx before uploading.

        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        logger.info("[FILE_PROCESSOR] Starting Word text extraction")
        metadata = {
            'extraction_method': 'python-docx',
        }

        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            text_parts = []
            paragraphs_count = 0
            tables_count = len(doc.tables)

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
                    paragraphs_count += 1

            # Also extract from tables
            for table in doc.tables:
                text_parts.append("\n[Tabel]")
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(' |'):
                        text_parts.append(row_text)

            extracted_text = '\n\n'.join(text_parts)
            metadata['paragraphs'] = paragraphs_count
            metadata['tables'] = tables_count

            if not extracted_text.strip():
                metadata['is_empty'] = True
                return "", metadata

            logger.info("[FILE_PROCESSOR] Word extraction complete: %d chars, %d paragraphs, %d tables",
                       len(extracted_text), paragraphs_count, tables_count)

            return extracted_text, metadata

        except ImportError:
            logger.error("[FILE_PROCESSOR] python-docx not installed")
            metadata['error'] = 'python-docx not installed'
            return "", metadata
        except ValueError as e:
            # Check if it's old .doc format
            error_str = str(e).lower()
            if 'not a word file' in error_str or 'content type' in error_str:
                logger.warning("[FILE_PROCESSOR] File is not a valid .docx (maybe old .doc format?)")
                metadata['error'] = 'Format .doc (Word 97-2003) tidak didukung. Silakan convert ke .docx (Word 2007+) dan upload ulang.'
                metadata['unsupported_format'] = 'doc_legacy'
                return "", metadata
            logger.exception("[FILE_PROCESSOR] Word extraction failed: %s", e)
            metadata['error'] = str(e)
            return "", metadata
        except Exception as e:
            logger.exception("[FILE_PROCESSOR] Word extraction failed: %s", e)
            metadata['error'] = str(e)
            return "", metadata

    @classmethod
    async def extract_excel_text(cls, content: bytes) -> Tuple[str, dict]:
        """
        Extract text from Excel spreadsheet as markdown tables.

        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        logger.info("[FILE_PROCESSOR] Starting Excel text extraction")
        metadata = {
            'extraction_method': 'pandas',
        }

        try:
            import pandas as pd

            excel_file = pd.ExcelFile(io.BytesIO(content))
            text_parts = []
            sheets_info = []

            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)

                # Skip empty sheets
                if df.empty:
                    continue

                text_parts.append(f"## Sheet: {sheet_name}\n")

                # Convert to markdown table
                if len(df.columns) > 0:
                    # Header
                    header = "| " + " | ".join(str(col) for col in df.columns) + " |"
                    separator = "| " + " | ".join("---" for _ in df.columns) + " |"
                    text_parts.append(header)
                    text_parts.append(separator)

                    # Rows (limit to first 100 rows to avoid huge output)
                    max_rows = 100
                    for _, row in df.head(max_rows).iterrows():
                        row_text = "| " + " | ".join(str(val) if pd.notna(val) else "" for val in row) + " |"
                        text_parts.append(row_text)

                    if len(df) > max_rows:
                        text_parts.append(f"\n... dan {len(df) - max_rows} baris lainnya")

                text_parts.append('')
                sheets_info.append({
                    'name': sheet_name,
                    'rows': len(df),
                    'columns': len(df.columns)
                })

            metadata['sheets'] = sheets_info
            metadata['total_sheets'] = len(sheets_info)

            extracted_text = '\n'.join(text_parts)

            if not extracted_text.strip():
                metadata['is_empty'] = True
                return "", metadata

            logger.info("[FILE_PROCESSOR] Excel extraction complete: %d chars, %d sheets",
                       len(extracted_text), len(sheets_info))

            return extracted_text, metadata

        except ImportError:
            logger.error("[FILE_PROCESSOR] pandas/openpyxl not installed")
            metadata['error'] = 'pandas/openpyxl not installed'
            return "", metadata
        except Exception as e:
            logger.exception("[FILE_PROCESSOR] Excel extraction failed: %s", e)
            metadata['error'] = str(e)
            return "", metadata

    @classmethod
    def _excel_to_images_sync(cls, content: bytes, max_rows_per_sheet: int = 50) -> Tuple[List[Image.Image], dict]:
        """Synchronous Excel to images conversion (run in thread pool).

        Converts each sheet to an image using matplotlib.

        Args:
            content: Excel file content as bytes
            max_rows_per_sheet: Maximum rows to render per sheet

        Returns:
            Tuple of (list of PIL Images, metadata dict)
        """
        import pandas as pd
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')  # Non-interactive backend

        metadata = {'sheets': [], 'total_sheets': 0}
        images = []

        excel_file = pd.ExcelFile(io.BytesIO(content))

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)

            if df.empty:
                continue

            # Limit rows
            df_display = df.head(max_rows_per_sheet)

            # Create figure
            fig, ax = plt.subplots(figsize=(12, min(20, 2 + len(df_display) * 0.3)))
            ax.axis('off')

            # Create table
            table = ax.table(
                cellText=df_display.values.tolist(),
                colLabels=[str(col) for col in df_display.columns],
                loc='upper center',
                cellLoc='center'
            )
            table.auto_set_font_size(False)
            table.set_fontsize(8)
            table.scale(1, 1.5)

            # Add sheet name as title
            plt.title(f"Sheet: {sheet_name}", fontsize=12, fontweight='bold', pad=20)

            # Convert to image
            buf = io.BytesIO()
            plt.savefig(buf, format='png', bbox_inches='tight', dpi=100, facecolor='white')
            buf.seek(0)
            img = Image.open(buf)
            images.append(img.copy())
            buf.close()
            plt.close(fig)

            metadata['sheets'].append({
                'name': sheet_name,
                'rows': len(df),
                'columns': len(df.columns),
                'displayed_rows': len(df_display)
            })

        metadata['total_sheets'] = len(images)

        return images, metadata

    @classmethod
    async def excel_to_images(cls, content: bytes, max_rows_per_sheet: int = 50, timeout: float = 60.0) -> Tuple[List[Image.Image], dict]:
        """
        Convert Excel sheets to images.

        Runs in a thread pool to avoid blocking the async event loop.

        Args:
            content: Excel file content as bytes
            max_rows_per_sheet: Maximum rows to render per sheet
            timeout: Maximum time in seconds for conversion

        Returns:
            Tuple of (list of PIL Images, metadata dict)
        """
        try:
            logger.info("[FILE_PROCESSOR] Starting Excel to image conversion")
            loop = asyncio.get_event_loop()

            images, metadata = await asyncio.wait_for(
                loop.run_in_executor(
                    _executor,
                    cls._excel_to_images_sync,
                    content,
                    max_rows_per_sheet
                ),
                timeout=timeout
            )

            logger.info("[FILE_PROCESSOR] Converted %d Excel sheets to images", len(images))
            return images, metadata

        except asyncio.TimeoutError:
            logger.error("[FILE_PROCESSOR] Excel to image conversion timed out")
            raise asyncio.TimeoutError("Excel conversion timed out")
        except ImportError as e:
            logger.error("[FILE_PROCESSOR] Missing dependency for Excel to image: %s", e)
            raise ImportError(f"Missing dependency: {e}")
        except Exception as e:
            logger.exception("[FILE_PROCESSOR] Excel to image conversion failed: %s", e)
            raise

    @classmethod
    async def extract_powerpoint_text(cls, content: bytes) -> Tuple[str, dict]:
        """
        Extract text from PowerPoint presentation.

        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        logger.info("[FILE_PROCESSOR] Starting PowerPoint text extraction")
        metadata = {
            'extraction_method': 'python-pptx',
        }

        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []
            slides_with_content = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"## Slide {slide_num}"]
                has_content = False

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        has_content = True

                if has_content:
                    text_parts.append('\n'.join(slide_text))
                    slides_with_content += 1

            metadata['total_slides'] = len(prs.slides)
            metadata['slides_with_content'] = slides_with_content

            extracted_text = '\n\n'.join(text_parts)

            if not extracted_text.strip():
                metadata['is_empty'] = True
                return "", metadata

            logger.info("[FILE_PROCESSOR] PowerPoint extraction complete: %d chars, %d slides",
                       len(extracted_text), slides_with_content)

            return extracted_text, metadata

        except ImportError:
            logger.error("[FILE_PROCESSOR] python-pptx not installed")
            metadata['error'] = 'python-pptx not installed'
            return "", metadata
        except Exception as e:
            logger.exception("[FILE_PROCESSOR] PowerPoint extraction failed: %s", e)
            metadata['error'] = str(e)
            return "", metadata

    @classmethod
    def is_vision_type(cls, mime_type: str) -> bool:
        """Check if file type should be processed with Vision AI."""
        return mime_type in cls.VISION_TYPES

    @classmethod
    def _pdf_to_images_sync(cls, content: bytes, dpi: int, poppler_path: Optional[str] = None) -> List[Image.Image]:
        """Synchronous PDF to image conversion (run in thread pool).

        Args:
            content: PDF bytes
            dpi: Resolution for rendering
            poppler_path: Path to poppler binaries (optional)

        Returns:
            List of PIL Images
        """
        logger.info("[FILE_PROCESSOR_SYNC] Starting PDF conversion (dpi=%d, size=%d bytes)", dpi, len(content))

        try:
            from pdf2image import convert_from_bytes

            # Build kwargs for convert_from_bytes
            kwargs = {
                "dpi": dpi,
                "fmt": "png",
                "thread_count": 2,  # Use 2 threads for faster conversion
            }

            # Add poppler path if available
            if poppler_path:
                kwargs["poppler_path"] = poppler_path
                logger.info("[FILE_PROCESSOR_SYNC] Using poppler path: %s", poppler_path)
            else:
                logger.warning("[FILE_PROCESSOR_SYNC] No poppler path specified, using system PATH")

            logger.info("[FILE_PROCESSOR_SYNC] Calling convert_from_bytes...")
            images = convert_from_bytes(content, **kwargs)
            logger.info("[FILE_PROCESSOR_SYNC] Conversion successful, got %d pages", len(images))

            return images

        except ImportError as e:
            logger.error("[FILE_PROCESSOR_SYNC] pdf2image not installed: %s", e)
            raise
        except Exception as e:
            logger.exception("[FILE_PROCESSOR_SYNC] PDF conversion failed: %s", e)
            raise

    @classmethod
    async def pdf_to_images(cls, content: bytes, dpi: int = 150, max_pages: int = 20, timeout: float = 120.0) -> List[Image.Image]:
        """
        Convert PDF pages to a list of PIL Images.

        Runs in a thread pool to avoid blocking the async event loop.
        Includes timeout to prevent indefinite hanging.

        Args:
            content: PDF file content as bytes
            dpi: Resolution for rendering (higher = better quality but larger)
            max_pages: Maximum pages to convert (to prevent memory issues)
            timeout: Maximum time in seconds for conversion (default 120s)

        Returns:
            List of PIL Image objects, one per page (up to max_pages)

        Raises:
            asyncio.TimeoutError: If conversion exceeds timeout
            ImportError: If pdf2image is not installed
            RuntimeError: If poppler is not installed or not working
        """
        # Verify poppler is installed first
        if not verify_poppler_installed():
            logger.error("[FILE_PROCESSOR] Poppler not installed or not working!")
            raise RuntimeError(
                "Poppler (pdftoppm) is not installed or not working. "
                "Install with: apt-get install poppler-utils (Linux) "
                "or brew install poppler (macOS)"
            )

        poppler_path = get_poppler_path()
        logger.info("[FILE_PROCESSOR] Starting PDF to image conversion at %d DPI (max %d pages, timeout %.1fs)",
                   dpi, max_pages, timeout)
        logger.info("[FILE_PROCESSOR] PDF content size: %d bytes (%.2f MB)", len(content), len(content) / (1024 * 1024))

        try:
            loop = asyncio.get_event_loop()

            # Run with timeout to prevent hanging
            logger.info("[FILE_PROCESSOR] Submitting to thread pool executor...")
            try:
                images = await asyncio.wait_for(
                    loop.run_in_executor(
                        _executor,
                        cls._pdf_to_images_sync,
                        content,
                        dpi,
                        poppler_path
                    ),
                    timeout=timeout
                )
            except asyncio.TimeoutError:
                logger.error("[FILE_PROCESSOR] PDF conversion timed out after %.1fs", timeout)
                raise asyncio.TimeoutError(f"PDF conversion timed out after {timeout} seconds")

            logger.info("[FILE_PROCESSOR] Thread pool returned %d images", len(images))

            # Limit pages to prevent memory issues
            if len(images) > max_pages:
                logger.warning("[FILE_PROCESSOR] PDF has %d pages, limiting to %d",
                             len(images), max_pages)
                images = images[:max_pages]

            logger.info("[FILE_PROCESSOR] Successfully converted %d pages", len(images))
            return images

        except ImportError:
            logger.error("[FILE_PROCESSOR] pdf2image not installed. Install with: pip install pdf2image")
            raise ImportError("pdf2image not installed. Install with: pip install pdf2image")
        except asyncio.TimeoutError:
            raise
        except Exception as e:
            logger.exception("[FILE_PROCESSOR] Failed to convert PDF to images: %s", e)
            raise

    @classmethod
    def concatenate_images_vertically(
        cls,
        images: List[Image.Image],
        max_height: int = 8000,
        padding: int = 20
    ) -> Image.Image:
        """
        Concatenate multiple images into one long vertical image.

        Args:
            images: List of PIL Image objects
            max_height: Maximum height in pixels (to avoid huge images)
            padding: Pixels between images

        Returns:
            Single concatenated PIL Image
        """
        if not images:
            raise ValueError("No images to concatenate")

        if len(images) == 1:
            return images[0]

        # Calculate total height
        total_height = sum(img.height for img in images) + padding * (len(images) - 1)

        # If too tall, scale down
        scale = 1.0
        if total_height > max_height:
            scale = max_height / total_height
            logger.info("[FILE_PROCESSOR] Scaling images by %.2f to fit max height %d", scale, max_height)

        # Find max width (after scaling)
        max_width = max(int(img.width * scale) for img in images)

        # Create new image
        final_height = int(total_height * scale)
        concatenated = Image.new('RGB', (max_width, final_height), color='white')

        y_offset = 0
        for img in images:
            if scale != 1.0:
                new_width = int(img.width * scale)
                new_height = int(img.height * scale)
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

            # Center horizontally if widths differ
            x_offset = (max_width - img.width) // 2
            concatenated.paste(img, (x_offset, y_offset))
            y_offset += img.height + padding

        logger.info("[FILE_PROCESSOR] Concatenated %d images into %dx%d",
                   len(images), max_width, final_height)
        return concatenated

    @classmethod
    def image_to_bytes(cls, image: Image.Image, format: str = 'PNG') -> bytes:
        """Convert PIL Image to bytes."""
        buffer = io.BytesIO()
        # Convert to RGB if necessary (for PNG compatibility)
        if image.mode in ('RGBA', 'P'):
            image = image.convert('RGB')
        image.save(buffer, format=format)
        return buffer.getvalue()


def get_mime_type(filename: str) -> str:
    """Get MIME type from filename extension."""
    import mimetypes
    mime_type, _ = mimetypes.guess_type(filename)
    return mime_type or 'application/octet-stream'
